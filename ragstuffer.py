#!/usr/bin/env python3
"""ragstuffer — ingests documents from Google Drive, git repos, and web URLs into Qdrant.

Document loading patterns (git shallow clone, glob filtering, web extraction, chunking
with source attribution) are adapted from the Red Hat Validated Patterns vector-embedder:
https://github.com/validatedpatterns-sandbox/vector-embedder
"""

import hashlib
import json
import logging
import os
import shutil
import subprocess
import sys
import time
import uuid
from datetime import UTC
from pathlib import Path

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    stream=sys.stdout,
)
log = logging.getLogger("ragstuffer")

# ── Configuration ────────────────────────────────────────────────────────────

STATE_PATH = Path(
    os.environ.get(
        "RAG_STATE_FILE",
        Path.home() / ".local" / "share" / "ramalama" / "rag-state.json",
    )
)

SUPPORTED_TEXT_EXTENSIONS = {".md", ".txt", ".adoc", ".rst"}
# HTML is text but needs tag stripping — handled separately in extract_text()
SUPPORTED_HTML_EXTENSIONS = {".html", ".htm"}
SUPPORTED_BINARY_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
ALL_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_HTML_EXTENSIONS | SUPPORTED_BINARY_EXTENSIONS


def load_state() -> dict:
    if STATE_PATH.exists():
        return json.loads(STATE_PATH.read_text())
    return {}


def save_state(state: dict) -> None:
    STATE_PATH.parent.mkdir(parents=True, exist_ok=True)
    STATE_PATH.write_text(json.dumps(state, indent=2))


# ── Text extraction ──────────────────────────────────────────────────────────


def extract_text(file_path: Path) -> str:
    """Extract text from a file based on its extension. Logs warnings on failure."""
    ext = file_path.suffix.lower()

    if ext in SUPPORTED_TEXT_EXTENSIONS:
        return file_path.read_text(errors="replace")

    if ext in SUPPORTED_HTML_EXTENSIONS:
        from html.parser import HTMLParser

        class _HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                self._skip = tag in ("script", "style", "nav", "header", "footer")

            def handle_endtag(self, tag):
                if tag in ("script", "style", "nav", "header", "footer"):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    t = data.strip()
                    if t:
                        self.parts.append(t)

        try:
            raw = file_path.read_text(errors="replace")
            parser = _HTMLTextExtractor()
            parser.feed(raw)
            return "\n".join(parser.parts)
        except Exception:
            log.warning("Failed to extract text from HTML: %s", file_path.name)
            return ""

    if ext == ".pdf":
        try:
            import pypdf

            reader = pypdf.PdfReader(str(file_path))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            log.warning("Failed to extract text from PDF: %s", file_path.name)
            return ""

    if ext == ".docx":
        try:
            import docx

            doc = docx.Document(str(file_path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            log.warning("Failed to extract text from DOCX: %s", file_path.name)
            return ""

    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n\n".join(texts)
        except Exception:
            log.warning("Failed to extract text from PPTX: %s", file_path.name)
            return ""

    if ext == ".xlsx":
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        texts.append(" | ".join(vals))
            return "\n".join(texts)
        except Exception:
            log.warning("Failed to extract text from XLSX: %s", file_path.name)
            return ""

    return ""


# ── Chunking (adapted from vector-embedder's RecursiveCharacterTextSplitter pattern) ─


def chunk_text(text: str, chunk_size: int = 1024, overlap: int = 128) -> list[str]:
    """Split text into overlapping chunks using paragraph/sentence boundaries.

    Tries to split on paragraph breaks first, then sentences, then words,
    falling back to character-level splits. Adapted from the chunking strategy
    in validatedpatterns-sandbox/vector-embedder.
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return [chunk for chunk in splitter.split_text(text) if chunk.strip()]


# ── Google Drive source ──────────────────────────────────────────────────────

GDRIVE_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

EXPORT_MAP = {
    "application/vnd.google-apps.document": ("application/pdf", ".pdf"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    ),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    ),
}


def get_drive_service():
    from google.oauth2 import service_account
    from googleapiclient.discovery import build

    creds_file = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
    if not creds_file:
        return None

    creds_path = Path(creds_file).expanduser()
    if not creds_path.exists():
        log.warning("Service account key not found: %s — Drive source disabled", creds_path)
        return None

    creds = service_account.Credentials.from_service_account_file(str(creds_path), scopes=GDRIVE_SCOPES)
    return build("drive", "v3", credentials=creds)


def list_drive_files(service, folder_id: str) -> list[dict]:
    results = []
    page_token = None
    while True:
        resp = (
            service.files()
            .list(
                q=f"'{folder_id}' in parents and trashed = false",
                fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
                pageSize=100,
                pageToken=page_token,
            )
            .execute()
        )
        results.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break
    return results


def download_drive_file(service, file_info: dict, staging_dir: Path) -> bool:
    from googleapiclient.http import MediaIoBaseDownload

    file_id = file_info["id"]
    name = file_info["name"]
    mime = file_info["mimeType"]

    if mime in EXPORT_MAP:
        export_mime, ext = EXPORT_MAP[mime]
        dest = staging_dir / f"{Path(name).stem}{ext}"
        log.info("Exporting %s → %s", name, dest.name)
        request = service.files().export_media(fileId=file_id, mimeType=export_mime)
    else:
        ext = Path(name).suffix.lower()
        if ext not in ALL_EXTENSIONS:
            log.debug("Skipping unsupported file: %s", name)
            return False
        dest = staging_dir / name
        log.info("Downloading %s", name)
        request = service.files().get_media(fileId=file_id)

    with open(dest, "wb") as fh:
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
    return True


def load_drive_docs(staging_dir: Path) -> list[dict]:
    """Load Drive documents: return list of {text, source} dicts."""
    folder_id = os.environ.get("GDRIVE_FOLDER_ID")
    if not folder_id:
        return []

    service = get_drive_service()
    if not service:
        return []

    state = load_state()
    files = list_drive_files(service, folder_id)
    changed = [f for f in files if state.get(f["id"]) != f["modifiedTime"]]

    if not changed:
        log.info("Drive: no new or modified files")
        return []

    log.info("Drive: %d new/modified file(s)", len(changed))
    # Prune state entries for files deleted from Drive
    current_ids = {f["id"] for f in files}
    new_state = {fid: mtime for fid, mtime in state.items() if fid in current_ids}

    for f in changed:
        try:
            if download_drive_file(service, f, staging_dir):
                new_state[f["id"]] = f["modifiedTime"]
        except Exception:
            log.warning("Drive: failed to download %s — skipping", f["name"])

    save_state(new_state)

    docs = []
    for f in staging_dir.iterdir():
        if not f.is_file():
            continue
        text = extract_text(f)
        if text:
            docs.append({"text": text, "source": f"gdrive://{f.name}"})
        else:
            log.warning("Drive: no text extracted from %s", f.name)
        # Clean up staging file after extraction
        f.unlink(missing_ok=True)
    return docs


# ── Git source (adapted from vector-embedder's GitLoader) ───────────────────


def load_git_docs(repos_dir: Path) -> list[dict]:
    """Clone/pull git repos and extract text from matching files.

    REPO_SOURCES is a JSON list of objects: [{"url": "...", "glob": "**/*.md"}, ...]
    Shallow clones and incremental pull adapted from vector-embedder.
    """
    sources_json = os.environ.get("REPO_SOURCES", "")
    if not sources_json:
        return []

    try:
        sources = json.loads(sources_json)
    except json.JSONDecodeError:
        log.error("REPO_SOURCES is not valid JSON")
        return []

    docs = []
    for source in sources:
        url = source.get("url", "")
        glob_pattern = source.get("glob", "**/*")
        if not url:
            continue

        repo_name = url.rstrip("/").split("/")[-1].replace(".git", "")
        repo_path = repos_dir / repo_name

        try:
            if repo_path.exists():
                log.info("Git: pulling %s", repo_name)
                result = subprocess.run(
                    ["git", "-C", str(repo_path), "pull", "--ff-only"],
                    capture_output=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    log.warning("Git: pull failed for %s, re-cloning", repo_name)
                    shutil.rmtree(repo_path)
                    subprocess.run(
                        ["git", "clone", "--depth", "1", url, str(repo_path)],
                        check=True,
                        capture_output=True,
                        timeout=120,
                    )
            else:
                log.info("Git: cloning %s (shallow)", repo_name)
                subprocess.run(
                    ["git", "clone", "--depth", "1", url, str(repo_path)],
                    check=True,
                    capture_output=True,
                    timeout=120,
                )
        except Exception:
            log.warning("Git: failed to clone/pull %s — skipping", url)
            continue

        for file_path in repo_path.glob(glob_pattern):
            if not file_path.is_file():
                continue
            if file_path.suffix.lower() not in ALL_EXTENSIONS:
                continue
            try:
                rel = file_path.relative_to(repo_path)
                text = extract_text(file_path)
                if text:
                    docs.append({"text": text, "source": f"{url}@{rel}"})
            except Exception:
                log.warning("Git: failed to extract %s — skipping", file_path)

    log.info("Git: loaded %d documents from %d repos", len(docs), len(sources))
    return docs


# ── Web source (adapted from vector-embedder's WebLoader) ───────────────────


def load_web_docs() -> list[dict]:
    """Fetch web pages and extract text.

    WEB_SOURCES is a JSON list of URLs: ["https://example.com/docs", ...]
    """
    sources_json = os.environ.get("WEB_SOURCES", "")
    if not sources_json:
        return []

    try:
        urls = json.loads(sources_json)
    except json.JSONDecodeError:
        log.error("WEB_SOURCES is not valid JSON")
        return []

    from html.parser import HTMLParser

    import requests

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            self._skip = tag in ("script", "style", "nav", "header", "footer")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "nav", "header", "footer"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip:
                text = data.strip()
                if text:
                    self.parts.append(text)

    docs = []
    for url in urls:
        try:
            log.info("Web: fetching %s", url)
            resp = requests.get(url, timeout=30, headers={"User-Agent": "ragstuffer/1.0"})
            resp.raise_for_status()

            if "html" in resp.headers.get("content-type", "").lower():
                parser = TextExtractor()
                parser.feed(resp.text)
                text = "\n\n".join(parser.parts)
            else:
                text = resp.text

            if text.strip():
                docs.append({"text": text, "source": url})
            else:
                log.warning("Web: no text extracted from %s", url)
        except Exception:
            log.warning("Web: failed to fetch %s — skipping", url)

    log.info("Web: loaded %d documents from %d URLs", len(docs), len(urls))
    return docs


# ── Qdrant ingestion ─────────────────────────────────────────────────────────


def get_qdrant_client():
    from qdrant_client import QdrantClient

    url = os.environ.get("QDRANT_URL", "http://127.0.0.1:6333")
    return QdrantClient(url=url, timeout=30)


def embed_texts(texts: list[str], batch_size: int = 64) -> list[list[float]]:
    """Embed texts via ragpipe's /v1/embeddings endpoint.

    Delegates to ragpipe instead of loading a local model — saves ~1-2 GB RAM.
    Batches requests to avoid oversized payloads.
    """
    import requests

    embed_url = os.environ.get("EMBED_URL", "http://127.0.0.1:8090/v1/embeddings")
    vectors = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i : i + batch_size]
        resp = requests.post(embed_url, json={"input": batch}, timeout=120)
        resp.raise_for_status()
        data = resp.json()["data"]
        # Sort by index to preserve order
        data.sort(key=lambda x: x["index"])
        vectors.extend([d["embedding"] for d in data])
    return vectors


def ensure_collection(qdrant, collection_name: str, vector_size: int):
    """Create Qdrant collection with scalar int8 quantization.

    Quantization reduces memory footprint of vectors. always_ram=True is
    required because HNSW rescoring needs quantized vectors in RAM for
    accurate distance computation during search.

    NOTE: Qdrant does not support adding quantization to an existing
    collection in place. If the collection already exists without
    quantization, it must be dropped and recreated.
    """
    from qdrant_client.models import Distance, ScalarQuantization, ScalarQuantizationConfig, ScalarType, VectorParams

    collections = [c.name for c in qdrant.get_collections().collections]
    if collection_name not in collections:
        qdrant.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
            quantization_config=ScalarQuantization(
                scalar=ScalarQuantizationConfig(
                    type=ScalarType.INT8,
                    quantile=0.99,
                    always_ram=True,
                ),
            ),
        )
        log.info("Created Qdrant collection '%s' with int8 scalar quantization", collection_name)


def _point_id(doc_id: str, chunk_id: int) -> int:
    """Deterministic point ID from doc_id:chunk_id for idempotent Qdrant upsert."""
    h = hashlib.sha256(f"{doc_id}:{chunk_id}".encode()).hexdigest()
    return int(h[:16], 16)


def _get_docstore():
    """Lazy-import and create docstore to avoid import at module level."""
    from docstore import create_docstore

    return create_docstore()


def ingest_docs(docs: list[dict]) -> bool:
    """Chunk documents, persist to docstore, embed, and upsert references to Qdrant.

    Qdrant payloads contain only {doc_id, chunk_id, source, created_at}.
    Full chunk text lives in the document store.
    """
    from datetime import datetime

    from qdrant_client.models import PointStruct

    if not docs:
        log.info("No documents to ingest")
        return True

    collection_name = os.environ.get("QDRANT_COLLECTION", "documents")
    chunk_size = int(os.environ.get("CHUNK_SIZE", "1024"))
    chunk_overlap = int(os.environ.get("CHUNK_OVERLAP", "128"))

    # Chunk all documents and assign stable doc_id per source
    all_chunks = []
    for doc in docs:
        doc_id = str(uuid.uuid5(uuid.NAMESPACE_URL, doc["source"]))
        chunks = chunk_text(doc["text"], chunk_size, chunk_overlap)
        for i, chunk in enumerate(chunks):
            all_chunks.append(
                {
                    "doc_id": doc_id,
                    "chunk_id": i,
                    "text": chunk,
                    "source": doc["source"],
                }
            )
        log.info("Chunked %s → %d chunks (doc_id=%s)", doc["source"], len(chunks), doc_id)

    if not all_chunks:
        log.warning("No text chunks to ingest")
        return True

    # Step 1: Persist to docstore (must succeed before Qdrant upsert)
    docstore = _get_docstore()
    docstore.upsert_chunks(all_chunks)
    log.info("Persisted %d chunks to docstore", len(all_chunks))

    # Step 2: Embed chunk texts via ragpipe (no local model needed)
    embed_batch_size = int(os.environ.get("EMBED_BATCH_SIZE", "64"))
    log.info("Embedding %d chunks via ragpipe (batch_size=%d)...", len(all_chunks), embed_batch_size)
    texts = [c["text"] for c in all_chunks]
    vectors = embed_texts(texts, batch_size=embed_batch_size)

    # Step 3: Upsert reference-only payloads to Qdrant
    qdrant = get_qdrant_client()

    # Ensure collection exists (idempotent — creates only if missing)
    ensure_collection(qdrant, collection_name, len(vectors[0]))

    now = datetime.now(UTC).isoformat()
    points = [
        PointStruct(
            id=_point_id(chunk["doc_id"], chunk["chunk_id"]),
            vector=vec,
            payload={
                "doc_id": chunk["doc_id"],
                "chunk_id": chunk["chunk_id"],
                "source": chunk["source"],
                "created_at": now,
            },
        )
        for vec, chunk in zip(vectors, all_chunks, strict=True)
    ]

    batch_size = 100
    for i in range(0, len(points), batch_size):
        qdrant.upsert(collection_name=collection_name, points=points[i : i + batch_size])

    log.info("Ingested %d reference points into Qdrant collection '%s'", len(points), collection_name)
    return True


# ── Poll loop ────────────────────────────────────────────────────────────────


def poll_once(staging_dir: Path, repos_dir: Path) -> None:
    """Single poll iteration: gather docs from all sources, ingest into Qdrant."""
    all_docs = []

    # Google Drive
    all_docs.extend(load_drive_docs(staging_dir))

    # Git repos
    all_docs.extend(load_git_docs(repos_dir))

    # Web URLs
    all_docs.extend(load_web_docs())

    if not all_docs:
        log.info("No documents to process from any source")
        return

    log.info("Total: %d documents from all sources", len(all_docs))

    if ingest_docs(all_docs):
        log.info("Qdrant updated — new documents available for RAG immediately")
    else:
        log.warning("Ingestion failed — will retry next poll")


def _start_admin_server(trigger_event, admin_port, admin_token):
    """Run a minimal HTTP admin server for triggering ingestion.

    Endpoints:
      POST /admin/ingest-now   — trigger immediate poll (incremental)
      POST /admin/ingest-full  — clear state + trigger full re-ingest
      GET  /health             — liveness check
    """
    import json as _json
    from http.server import BaseHTTPRequestHandler, HTTPServer
    import threading

    class AdminHandler(BaseHTTPRequestHandler):
        def _check_auth(self):
            if not admin_token:
                return True
            auth = self.headers.get("Authorization", "")
            if auth == f"Bearer {admin_token}":
                return True
            self.send_response(401)
            self.end_headers()
            self.wfile.write(b'{"error":"unauthorized"}')
            return False

        def do_GET(self):
            if self.path == "/health":
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(b'{"status":"ok"}')
            else:
                self.send_response(404)
                self.end_headers()

        def do_POST(self):
            if not self._check_auth():
                return

            if self.path == "/admin/ingest-now":
                log.info("Admin trigger: ingest-now (incremental)")
                trigger_event.set()
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(_json.dumps({"status": "triggered", "mode": "incremental"}).encode())

            elif self.path == "/admin/ingest-full":
                log.info("Admin trigger: ingest-full (clearing state)")
                if STATE_PATH.exists():
                    STATE_PATH.unlink()
                    log.info("Cleared state file: %s", STATE_PATH)
                trigger_event.set()
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(_json.dumps({"status": "triggered", "mode": "full"}).encode())

            else:
                self.send_response(404)
                self.end_headers()

        def log_message(self, format, *args):
            """Suppress default access logging — we log triggers explicitly."""

    server = HTTPServer(("0.0.0.0", admin_port), AdminHandler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    log.info("Admin server listening on :%d", admin_port)
    return server


def main() -> None:
    import threading

    interval = int(os.environ.get("WATCH_INTERVAL_MINUTES", "15"))
    admin_port = int(os.environ.get("RAGSTUFFER_ADMIN_PORT", "8091"))
    admin_token = os.environ.get("RAGSTUFFER_ADMIN_TOKEN", "")
    staging_dir = Path(os.environ.get("STAGING_DIR", "/tmp/rag-staging"))
    repos_dir = Path(os.environ.get("REPOS_DIR", "/tmp/rag-repos"))

    staging_dir.mkdir(parents=True, exist_ok=True)
    repos_dir.mkdir(parents=True, exist_ok=True)

    # Event for admin-triggered immediate ingestion
    trigger = threading.Event()

    log.info("Starting ragstuffer")
    log.info("  Qdrant:     %s", os.environ.get("QDRANT_URL", "http://127.0.0.1:6333"))
    log.info("  Collection: %s", os.environ.get("QDRANT_COLLECTION", "documents"))
    log.info("  Interval:   %d min", interval)
    log.info("  Admin:      :%d", admin_port)
    if os.environ.get("GDRIVE_FOLDER_ID"):
        log.info("  Drive:      folder %s", os.environ["GDRIVE_FOLDER_ID"])
    if os.environ.get("REPO_SOURCES"):
        log.info("  Git:        %s", os.environ["REPO_SOURCES"])
    if os.environ.get("WEB_SOURCES"):
        log.info("  Web:        %s", os.environ["WEB_SOURCES"])

    _start_admin_server(trigger, admin_port, admin_token)

    while True:
        try:
            poll_once(staging_dir, repos_dir)
        except Exception:
            log.exception("Poll cycle failed")
        # Wait for interval OR admin trigger, whichever comes first
        trigger.wait(timeout=interval * 60)
        trigger.clear()


if __name__ == "__main__":
    main()

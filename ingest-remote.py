#!/usr/bin/env python3
"""One-shot ingestion — runs on a remote host with FastEmbed locally.

Embeds documents using all available CPU cores, then pushes chunks to
Postgres and vectors to Qdrant on harrison (192.168.1.122).

Usage:
    # From harrison, deploy and run on lennon:
    # See deploy-to-lennon.sh

    # Or manually on lennon:
    python3 ingest-remote.py
"""

import hashlib
import json
import logging
import os
import shutil
import subprocess
import sys
import time
import uuid
from datetime import UTC, datetime
from pathlib import Path

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    stream=sys.stdout,
)
log = logging.getLogger("ingest-remote")

# ── Configuration ────────────────────────────────────────────────────────────

HARRISON = os.environ.get("HARRISON_HOST", "192.168.1.122")
QDRANT_URL = os.environ.get("QDRANT_URL", f"http://{HARRISON}:6333")
DOCSTORE_URL = os.environ.get("DOCSTORE_URL", f"postgresql://litellm:litellm@{HARRISON}:5432/litellm")
QDRANT_COLLECTION = os.environ.get("QDRANT_COLLECTION", "documents")
EMBED_MODEL = os.environ.get("EMBED_MODEL", "BAAI/bge-base-en-v1.5")
EMBED_THREADS = int(os.environ.get("EMBED_THREADS", os.cpu_count() or 4))
CHUNK_SIZE = int(os.environ.get("CHUNK_SIZE", "1024"))
CHUNK_OVERLAP = int(os.environ.get("CHUNK_OVERLAP", "128"))

# Google Drive
GDRIVE_FOLDER_ID = os.environ.get("GDRIVE_FOLDER_ID", "")
GDRIVE_SA_KEY = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS", str(Path.home() / ".config/ramalama/gdrive-sa.json"))

# Git repos: JSON list of {"url": "...", "glob": "**/*.md"}
REPO_SOURCES = os.environ.get("REPO_SOURCES", "")

# Web URLs: JSON list of URLs
WEB_SOURCES = os.environ.get("WEB_SOURCES", "")

SUPPORTED_TEXT_EXTENSIONS = {".md", ".txt", ".adoc", ".rst"}
SUPPORTED_HTML_EXTENSIONS = {".html", ".htm"}
SUPPORTED_BINARY_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
ALL_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_HTML_EXTENSIONS | SUPPORTED_BINARY_EXTENSIONS

GDRIVE_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
EXPORT_MAP = {
    "application/vnd.google-apps.document": ("application/pdf", ".pdf"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    ),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    ),
}


# ── Text extraction (same as rag-watcher.py) ────────────────────────────────


def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()

    if ext in SUPPORTED_TEXT_EXTENSIONS:
        return file_path.read_text(errors="replace")

    if ext in SUPPORTED_HTML_EXTENSIONS:
        from html.parser import HTMLParser

        class _HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                self._skip = tag in ("script", "style", "nav", "header", "footer")

            def handle_endtag(self, tag):
                if tag in ("script", "style", "nav", "header", "footer"):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    t = data.strip()
                    if t:
                        self.parts.append(t)

        try:
            raw = file_path.read_text(errors="replace")
            parser = _HTMLTextExtractor()
            parser.feed(raw)
            return "\n".join(parser.parts)
        except Exception:
            log.warning("Failed to extract text from HTML: %s", file_path.name)
            return ""

    if ext == ".pdf":
        try:
            import pypdf

            reader = pypdf.PdfReader(str(file_path))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            log.warning("Failed to extract text from PDF: %s", file_path.name)
            return ""

    if ext == ".docx":
        try:
            import docx

            doc = docx.Document(str(file_path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            log.warning("Failed to extract text from DOCX: %s", file_path.name)
            return ""

    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n\n".join(texts)
        except Exception:
            log.warning("Failed to extract text from PPTX: %s", file_path.name)
            return ""

    if ext == ".xlsx":
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        texts.append(" | ".join(vals))
            return "\n".join(texts)
        except Exception:
            log.warning("Failed to extract text from XLSX: %s", file_path.name)
            return ""

    return ""


# ── Chunking ─────────────────────────────────────────────────────────────────


def chunk_text(text: str) -> list[str]:
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return [chunk for chunk in splitter.split_text(text) if chunk.strip()]


# ── Google Drive ─────────────────────────────────────────────────────────────


def load_drive_docs(staging_dir: Path) -> list[dict]:
    if not GDRIVE_FOLDER_ID:
        log.info("Drive: no GDRIVE_FOLDER_ID set, skipping")
        return []

    sa_path = Path(GDRIVE_SA_KEY).expanduser()
    if not sa_path.exists():
        log.warning("Drive: SA key not found at %s, skipping", sa_path)
        return []

    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = str(sa_path)

    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload

    creds = service_account.Credentials.from_service_account_file(str(sa_path), scopes=GDRIVE_SCOPES)
    service = build("drive", "v3", credentials=creds)

    # List all files in folder
    files = []
    page_token = None
    while True:
        resp = (
            service.files()
            .list(
                q=f"'{GDRIVE_FOLDER_ID}' in parents and trashed = false",
                fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
                pageSize=100,
                pageToken=page_token,
            )
            .execute()
        )
        files.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    log.info("Drive: found %d files", len(files))

    for f in files:
        file_id = f["id"]
        name = f["name"]
        mime = f["mimeType"]

        if mime in EXPORT_MAP:
            export_mime, ext = EXPORT_MAP[mime]
            dest = staging_dir / f"{Path(name).stem}{ext}"
            request = service.files().export_media(fileId=file_id, mimeType=export_mime)
        else:
            ext = Path(name).suffix.lower()
            if ext not in ALL_EXTENSIONS:
                log.debug("Skipping unsupported: %s", name)
                continue
            dest = staging_dir / name
            request = service.files().get_media(fileId=file_id)

        log.info("Downloading %s", name)
        with open(dest, "wb") as fh:
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()

    docs = []
    for f in staging_dir.iterdir():
        if not f.is_file():
            continue
        text = extract_text(f)
        if text:
            docs.append({"text": text, "source": f"gdrive://{f.name}"})
        else:
            log.warning("Drive: no text from %s", f.name)
        f.unlink(missing_ok=True)

    return docs


# ── Git source ───────────────────────────────────────────────────────────────


def load_git_docs(repos_dir: Path) -> list[dict]:
    if not REPO_SOURCES:
        return []

    try:
        sources = json.loads(REPO_SOURCES)
    except json.JSONDecodeError:
        log.error("REPO_SOURCES is not valid JSON")
        return []

    docs = []
    for source in sources:
        url = source.get("url", "")
        glob_pattern = source.get("glob", "**/*")
        if not url:
            continue

        repo_name = url.rstrip("/").split("/")[-1].replace(".git", "")
        repo_path = repos_dir / repo_name

        try:
            if repo_path.exists():
                log.info("Git: pulling %s", repo_name)
                result = subprocess.run(
                    ["git", "-C", str(repo_path), "pull", "--ff-only"],
                    capture_output=True,
                    timeout=60,
                )
                if result.returncode != 0:
                    shutil.rmtree(repo_path)
                    subprocess.run(
                        ["git", "clone", "--depth", "1", url, str(repo_path)],
                        check=True,
                        capture_output=True,
                        timeout=120,
                    )
            else:
                log.info("Git: cloning %s (shallow)", repo_name)
                subprocess.run(
                    ["git", "clone", "--depth", "1", url, str(repo_path)],
                    check=True,
                    capture_output=True,
                    timeout=120,
                )
        except Exception:
            log.warning("Git: failed to clone/pull %s, skipping", url)
            continue

        for file_path in repo_path.glob(glob_pattern):
            if not file_path.is_file():
                continue
            if file_path.suffix.lower() not in ALL_EXTENSIONS:
                continue
            try:
                rel = file_path.relative_to(repo_path)
                text = extract_text(file_path)
                if text:
                    docs.append({"text": text, "source": f"{url}@{rel}"})
            except Exception:
                log.warning("Git: failed to extract %s", file_path)

    log.info("Git: loaded %d documents from %d repos", len(docs), len(sources))
    return docs


# ── Web source ───────────────────────────────────────────────────────────────


def load_web_docs() -> list[dict]:
    if not WEB_SOURCES:
        return []

    try:
        urls = json.loads(WEB_SOURCES)
    except json.JSONDecodeError:
        log.error("WEB_SOURCES is not valid JSON")
        return []

    from html.parser import HTMLParser

    import requests

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            self._skip = tag in ("script", "style", "nav", "header", "footer")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "nav", "header", "footer"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip:
                text = data.strip()
                if text:
                    self.parts.append(text)

    docs = []
    for url in urls:
        try:
            log.info("Web: fetching %s", url)
            resp = requests.get(url, timeout=30, headers={"User-Agent": "rag-watcher/1.0"})
            resp.raise_for_status()
            if "html" in resp.headers.get("content-type", "").lower():
                parser = TextExtractor()
                parser.feed(resp.text)
                text = "\n\n".join(parser.parts)
            else:
                text = resp.text
            if text.strip():
                docs.append({"text": text, "source": url})
        except Exception:
            log.warning("Web: failed to fetch %s", url)

    log.info("Web: loaded %d documents from %d URLs", len(docs), len(urls))
    return docs


# ── Local FastEmbed ──────────────────────────────────────────────────────────


def embed_texts_local(texts: list[str], batch_size: int = 256) -> list[list[float]]:
    """Embed using sentence-transformers on CUDA (4070 Ti)."""
    import torch
    from sentence_transformers import SentenceTransformer

    device = "cuda" if torch.cuda.is_available() else "cpu"
    log.info("Loading model %s on %s...", EMBED_MODEL, device)
    model = SentenceTransformer(EMBED_MODEL, device=device)

    log.info("Embedding %d texts (batch_size=%d, device=%s)...", len(texts), batch_size, device)
    t0 = time.monotonic()
    vectors = model.encode(
        texts,
        batch_size=batch_size,
        show_progress_bar=True,
        normalize_embeddings=True,
    )
    elapsed = time.monotonic() - t0
    log.info(
        "Embedded %d texts in %.1fs (%.0f texts/sec)",
        len(texts),
        elapsed,
        len(texts) / elapsed if elapsed > 0 else 0,
    )
    return vectors.tolist()


# ── Qdrant + docstore ───────────────────────────────────────────────────────


def _point_id(doc_id: str, chunk_id: int) -> int:
    h = hashlib.sha256(f"{doc_id}:{chunk_id}".encode()).hexdigest()
    return int(h[:16], 16)


def ingest(docs: list[dict]) -> None:
    import psycopg2
    from psycopg2.extras import execute_values
    from qdrant_client import QdrantClient
    from qdrant_client.models import (
        Distance,
        PointStruct,
        ScalarQuantization,
        ScalarQuantizationConfig,
        ScalarType,
        VectorParams,
    )

    if not docs:
        log.info("No documents to ingest")
        return

    # Chunk all documents
    all_chunks = []
    for doc in docs:
        doc_id = str(uuid.uuid5(uuid.NAMESPACE_URL, doc["source"]))
        chunks = chunk_text(doc["text"])
        for i, chunk in enumerate(chunks):
            all_chunks.append(
                {
                    "doc_id": doc_id,
                    "chunk_id": i,
                    "text": chunk,
                    "source": doc["source"],
                }
            )
        log.info("Chunked %s -> %d chunks (doc_id=%s)", doc["source"], len(chunks), doc_id)

    if not all_chunks:
        log.warning("No text chunks produced")
        return

    log.info("Total: %d chunks from %d documents", len(all_chunks), len(docs))

    # Step 1: Embed locally with FastEmbed
    texts = [c["text"] for c in all_chunks]
    vectors = embed_texts_local(texts)

    # Step 2: Push chunks to Postgres docstore on harrison
    log.info("Connecting to Postgres at %s...", HARRISON)
    conn = psycopg2.connect(DOCSTORE_URL)
    conn.autocommit = True
    with conn.cursor() as cur:
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chunks (
                doc_id     TEXT NOT NULL,
                chunk_id   INTEGER NOT NULL,
                text       TEXT NOT NULL,
                source     TEXT NOT NULL DEFAULT '',
                created_at TEXT NOT NULL DEFAULT '',
                PRIMARY KEY (doc_id, chunk_id)
            )
        """)
        now = datetime.now(UTC).isoformat()
        values = [(c["doc_id"], c["chunk_id"], c["text"], c["source"], now) for c in all_chunks]
        execute_values(
            cur,
            """
            INSERT INTO chunks (doc_id, chunk_id, text, source, created_at)
            VALUES %s
            ON CONFLICT (doc_id, chunk_id)
            DO UPDATE SET text = EXCLUDED.text, source = EXCLUDED.source
            """,
            values,
        )
    conn.close()
    log.info("Persisted %d chunks to Postgres", len(all_chunks))

    # Step 3: Upsert vectors to Qdrant on harrison
    log.info("Connecting to Qdrant at %s...", QDRANT_URL)
    qdrant = QdrantClient(url=QDRANT_URL, timeout=60)

    collections = [c.name for c in qdrant.get_collections().collections]
    if QDRANT_COLLECTION not in collections:
        qdrant.create_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=VectorParams(size=len(vectors[0]), distance=Distance.COSINE),
            quantization_config=ScalarQuantization(
                scalar=ScalarQuantizationConfig(
                    type=ScalarType.INT8,
                    quantile=0.99,
                    always_ram=True,
                ),
            ),
        )
        log.info("Created Qdrant collection '%s'", QDRANT_COLLECTION)

    now = datetime.now(UTC).isoformat()
    points = [
        PointStruct(
            id=_point_id(chunk["doc_id"], chunk["chunk_id"]),
            vector=vec,
            payload={
                "doc_id": chunk["doc_id"],
                "chunk_id": chunk["chunk_id"],
                "source": chunk["source"],
                "created_at": now,
            },
        )
        for vec, chunk in zip(vectors, all_chunks, strict=True)
    ]

    batch_size = 100
    for i in range(0, len(points), batch_size):
        qdrant.upsert(collection_name=QDRANT_COLLECTION, points=points[i : i + batch_size])
        if (i // batch_size) % 10 == 0:
            log.info("Qdrant upsert: %d/%d points", min(i + batch_size, len(points)), len(points))

    log.info("Ingested %d points into Qdrant collection '%s'", len(points), QDRANT_COLLECTION)


# ── Main ─────────────────────────────────────────────────────────────────────


def main() -> None:
    staging_dir = Path("/tmp/rag-staging")
    repos_dir = Path("/tmp/rag-repos")
    staging_dir.mkdir(parents=True, exist_ok=True)
    repos_dir.mkdir(parents=True, exist_ok=True)

    log.info("One-shot ingestion — embedding locally, pushing to %s", HARRISON)
    log.info("  Qdrant:     %s", QDRANT_URL)
    log.info("  Postgres:   %s@%s", "litellm", HARRISON)
    log.info("  Collection: %s", QDRANT_COLLECTION)
    log.info("  Model:      %s", EMBED_MODEL)
    log.info("  Threads:    %d", EMBED_THREADS)

    all_docs = []
    all_docs.extend(load_drive_docs(staging_dir))
    all_docs.extend(load_git_docs(repos_dir))
    all_docs.extend(load_web_docs())

    if not all_docs:
        log.info("No documents found from any source")
        return

    log.info("Total: %d documents from all sources", len(all_docs))

    t0 = time.monotonic()
    ingest(all_docs)
    elapsed = time.monotonic() - t0
    log.info("Done in %.1fs", elapsed)


if __name__ == "__main__":
    main()

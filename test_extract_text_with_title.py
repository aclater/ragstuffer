"""Tests for extract_text_with_title() — title extraction and fallback chain.

Covers all source types with real file fixtures:
  - PDF with/without /Title metadata
  - DOCX with/without core_properties.title
  - PPTX with/without core_properties.title
  - XLSX (always falls back to stem)
  - HTML with/without <title> tag
  - Markdown with/without # heading
  - RST with/without underline title
  - AsciiDoc with/without = heading
  - Plain text (always falls back to stem)
"""

import pypdf
import docx
from pptx import Presentation
import openpyxl

from common import ExtractedText, extract_text_with_title


# ── Text formats ────────────────────────────────────────────────────────────


class TestMarkdown:
    def test_with_heading(self, tmp_path):
        f = tmp_path / "readme.md"
        f.write_text("# My Project\n\nSome content here.")
        result = extract_text_with_title(f)
        assert result.title == "My Project"
        assert "Some content" in result.text

    def test_without_heading(self, tmp_path):
        f = tmp_path / "notes.md"
        f.write_text("Just some notes without a heading.")
        result = extract_text_with_title(f)
        assert result.title == "notes"  # stem fallback
        assert "Just some notes" in result.text


class TestRst:
    def test_with_underline_title(self, tmp_path):
        f = tmp_path / "guide.rst"
        f.write_text("User Guide\n==========\n\nWelcome.")
        result = extract_text_with_title(f)
        assert result.title == "User Guide"
        assert "Welcome" in result.text

    def test_without_title(self, tmp_path):
        f = tmp_path / "fragment.rst"
        f.write_text("Just a paragraph of text.")
        result = extract_text_with_title(f)
        assert result.title == "fragment"  # stem fallback


class TestAsciidoc:
    def test_with_heading(self, tmp_path):
        f = tmp_path / "manual.adoc"
        f.write_text("= Operations Manual\n\nFirst section.")
        result = extract_text_with_title(f)
        assert result.title == "Operations Manual"
        assert "First section" in result.text

    def test_without_heading(self, tmp_path):
        f = tmp_path / "snippet.adoc"
        f.write_text("Some AsciiDoc content here.")
        result = extract_text_with_title(f)
        assert result.title == "snippet"  # stem fallback


class TestPlainText:
    def test_falls_back_to_stem(self, tmp_path):
        f = tmp_path / "changelog.txt"
        f.write_text("v1.0.0 — initial release")
        result = extract_text_with_title(f)
        assert result.title == "changelog"  # always stem
        assert "v1.0.0" in result.text

    def test_empty_file(self, tmp_path):
        f = tmp_path / "empty.txt"
        f.write_text("")
        result = extract_text_with_title(f)
        assert result.text == ""
        # Empty text files still get stem as title (read succeeds, text is empty)
        assert result.title == "empty"


# ── HTML ────────────────────────────────────────────────────────────────────


class TestHtml:
    def test_with_title_tag(self, tmp_path):
        f = tmp_path / "page.html"
        f.write_text("<html><head><title>My Page</title></head><body><p>Content</p></body></html>")
        result = extract_text_with_title(f)
        assert result.title == "My Page"
        assert "Content" in result.text

    def test_without_title_tag(self, tmp_path):
        f = tmp_path / "notitle.html"
        f.write_text("<html><body><p>No title here</p></body></html>")
        result = extract_text_with_title(f)
        assert result.title == "notitle"  # stem fallback
        assert "No title here" in result.text

    def test_htm_extension(self, tmp_path):
        f = tmp_path / "legacy.htm"
        f.write_text("<html><head><title>Legacy</title></head><body><p>Old page</p></body></html>")
        result = extract_text_with_title(f)
        assert result.title == "Legacy"


# ── PDF ─────────────────────────────────────────────────────────────────────


class TestPdf:
    def test_with_metadata_title(self, tmp_path):
        f = tmp_path / "titled.pdf"
        writer = pypdf.PdfWriter()
        writer.add_blank_page(width=72, height=72)
        writer.add_metadata({"/Title": "PDF Document Title"})
        with open(f, "wb") as fh:
            writer.write(fh)

        result = extract_text_with_title(f)
        assert result.title == "PDF Document Title"

    def test_without_metadata_title(self, tmp_path):
        f = tmp_path / "untitled.pdf"
        writer = pypdf.PdfWriter()
        writer.add_blank_page(width=72, height=72)
        with open(f, "wb") as fh:
            writer.write(fh)

        result = extract_text_with_title(f)
        assert result.title == "untitled"  # stem fallback


# ── DOCX ────────────────────────────────────────────────────────────────────


class TestDocx:
    def test_with_title(self, tmp_path):
        f = tmp_path / "report.docx"
        doc = docx.Document()
        doc.core_properties.title = "Quarterly Report"
        doc.add_paragraph("Revenue increased 15%.")
        doc.save(str(f))

        result = extract_text_with_title(f)
        assert result.title == "Quarterly Report"
        assert "Revenue" in result.text

    def test_without_title(self, tmp_path):
        f = tmp_path / "draft.docx"
        doc = docx.Document()
        doc.add_paragraph("Draft content.")
        doc.save(str(f))

        result = extract_text_with_title(f)
        assert result.title == "draft"  # stem fallback
        assert "Draft content" in result.text


# ── PPTX ────────────────────────────────────────────────────────────────────


class TestPptx:
    def test_with_title(self, tmp_path):
        f = tmp_path / "deck.pptx"
        prs = Presentation()
        prs.core_properties.title = "Strategy Deck"
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        prs.save(str(f))

        result = extract_text_with_title(f)
        assert result.title == "Strategy Deck"

    def test_without_title(self, tmp_path):
        f = tmp_path / "slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(f))

        result = extract_text_with_title(f)
        assert result.title == "slides"  # stem fallback


# ── XLSX ────────────────────────────────────────────────────────────────────


class TestXlsx:
    def test_always_uses_stem(self, tmp_path):
        f = tmp_path / "budget.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Item", "Cost"])
        ws.append(["Widget", 42])
        wb.save(str(f))

        result = extract_text_with_title(f)
        assert result.title == "budget"  # always stem
        assert "Widget" in result.text
        assert "42" in result.text


# ── Return type ─────────────────────────────────────────────────────────────


class TestReturnType:
    def test_returns_extracted_text(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("hello")
        result = extract_text_with_title(f)
        assert isinstance(result, ExtractedText)

    def test_unsupported_extension_returns_empty(self, tmp_path):
        f = tmp_path / "data.bin"
        f.write_bytes(b"\x00\x01\x02")
        result = extract_text_with_title(f)
        assert result.text == ""
        assert result.title == ""

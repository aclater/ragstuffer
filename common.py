"""Shared constants and utilities for ragstuffer and ingest-remote.

Consolidates text extraction, chunking, file type constants, and Google Drive
constants that were previously duplicated between the two scripts.
"""

import logging
import re
from dataclasses import dataclass
from pathlib import Path

log = logging.getLogger("ragstuffer.common")

# ── Lazy-loaded optional dependencies ─────────────────────────────────────────
# Cached at module level to avoid repeated sys.modules lookups across
# hundreds of extract_text() calls during bulk ingestion.

_pypdf = None
_docx = None
_Presentation = None
_openpyxl = None


def _get_pypdf():
    global _pypdf
    if _pypdf is None:
        import pypdf
        _pypdf = pypdf
    return _pypdf


def _get_docx():
    global _docx
    if _docx is None:
        import docx
        _docx = docx
    return _docx


def _get_presentation():
    global _Presentation
    if _Presentation is None:
        from pptx import Presentation
        _Presentation = Presentation
    return _Presentation


def _get_openpyxl():
    global _openpyxl
    if _openpyxl is None:
        import openpyxl
        _openpyxl = openpyxl
    return _openpyxl


# ── File type constants ───────────────────────────────────────────────────────

SUPPORTED_TEXT_EXTENSIONS = {".md", ".txt", ".adoc", ".rst"}
SUPPORTED_HTML_EXTENSIONS = {".html", ".htm"}
SUPPORTED_BINARY_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
ALL_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_HTML_EXTENSIONS | SUPPORTED_BINARY_EXTENSIONS

# ── Google Drive constants ────────────────────────────────────────────────────

GDRIVE_SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

EXPORT_MAP = {
    "application/vnd.google-apps.document": ("application/pdf", ".pdf"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    ),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    ),
}


# ── ExtractedText ────────────────────────────────────────────────────────────


@dataclass
class ExtractedText:
    text: str
    title: str = ""


# ── HTML text extraction ─────────────────────────────────────────────────────

_TITLE_RE = re.compile(r"<title[^>]*>(.*?)</title>", re.DOTALL | re.IGNORECASE)


def _make_html_parser():
    """Create an HTML parser that strips scripts/styles/nav and returns text."""
    from html.parser import HTMLParser

    class _HTMLTextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts: list[str] = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            self._skip = tag in ("script", "style", "nav", "header", "footer")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "nav", "header", "footer"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip:
                t = data.strip()
                if t:
                    self.parts.append(t)

    return _HTMLTextExtractor()


def _extract_html_title(raw: str) -> str:
    """Extract the <title> tag content from raw HTML, or return empty string."""
    m = _TITLE_RE.search(raw)
    if m:
        return m.group(1).strip()
    return ""


# ── Text extraction ──────────────────────────────────────────────────────────


def _extract_title_from_text(content: str) -> str:
    """Extract a title from plain-text content (Markdown heading, AsciiDoc, RST)."""
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith("# ") and not stripped.startswith("# " * 2):
            return stripped[2:].strip()
        if stripped.startswith("= ") and not stripped.startswith("== "):
            return stripped[2:].strip()
        break
    return ""


def _extract_title_from_rst(content: str) -> str:
    """Extract title from reStructuredText (Title\n===== pattern)."""
    lines = content.splitlines()
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped and i + 1 < len(lines):
            next_line = lines[i + 1].strip()
            if next_line and all(c in ("=", "-", "~", "^") for c in next_line) and len(next_line) >= len(stripped):
                return stripped
    return ""


def extract_text(file_path: Path) -> str:
    """Extract text from a file based on its extension. Logs warnings on failure."""
    return extract_text_with_title(file_path).text


def extract_text_with_title(file_path: Path) -> ExtractedText:
    """Extract text and title from a file based on its extension.

    Title fallback chain per source type:
      - PDF:     /Title metadata → filename stem → ""
      - DOCX:    core_properties.title → filename stem → ""
      - PPTX:    core_properties.title → filename stem → ""
      - XLSX:    filename stem → ""
      - HTML:    <title> tag → filename stem → ""
      - Markdown: first # heading → filename stem → ""
      - RST:     Title\n===== → filename stem → ""
      - Adoc:    = Title → filename stem → ""
      - Plain text: filename stem → ""
    """
    ext = file_path.suffix.lower()
    stem = file_path.stem

    if ext in SUPPORTED_TEXT_EXTENSIONS:
        try:
            content = file_path.read_text(errors="replace")
        except Exception:
            log.warning("Failed to read text file: %s", file_path.name)
            return ExtractedText(text="", title="")

        title = ""
        if ext == ".md":
            title = _extract_title_from_text(content)
        elif ext == ".adoc":
            title = _extract_title_from_text(content)
        elif ext == ".rst":
            title = _extract_title_from_rst(content)

        return ExtractedText(text=content, title=title or stem)

    if ext in SUPPORTED_HTML_EXTENSIONS:
        try:
            raw = file_path.read_text(errors="replace")
            title = _extract_html_title(raw) or stem
            parser = _make_html_parser()
            parser.feed(raw)
            return ExtractedText(text="\n".join(parser.parts), title=title)
        except Exception:
            log.warning("Failed to extract text from HTML: %s", file_path.name)
            return ExtractedText(text="", title=stem)

    if ext == ".pdf":
        try:
            reader = _get_pypdf().PdfReader(str(file_path))
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            title = ""
            if reader.metadata and reader.metadata.get("/Title"):
                title = reader.metadata["/Title"].strip()
            return ExtractedText(text=text, title=title or stem)
        except Exception:
            log.warning("Failed to extract text from PDF: %s", file_path.name)
            return ExtractedText(text="", title=stem)

    if ext == ".docx":
        try:
            doc = _get_docx().Document(str(file_path))
            text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
            title = ""
            if doc.core_properties.title:
                title = doc.core_properties.title.strip()
            return ExtractedText(text=text, title=title or stem)
        except Exception:
            log.warning("Failed to extract text from DOCX: %s", file_path.name)
            return ExtractedText(text="", title=stem)

    if ext == ".pptx":
        try:
            prs = _get_presentation()(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            text = "\n\n".join(texts)
            title = ""
            if prs.core_properties.title:
                title = prs.core_properties.title.strip()
            return ExtractedText(text=text, title=title or stem)
        except Exception:
            log.warning("Failed to extract text from PPTX: %s", file_path.name)
            return ExtractedText(text="", title=stem)

    if ext == ".xlsx":
        try:
            wb = _get_openpyxl().load_workbook(str(file_path), read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None]
                    if vals:
                        texts.append(" | ".join(vals))
            text = "\n".join(texts)
            return ExtractedText(text=text, title=stem)
        except Exception:
            log.warning("Failed to extract text from XLSX: %s", file_path.name)
            return ExtractedText(text="", title=stem)

    return ExtractedText(text="", title="")


# ── Chunking ─────────────────────────────────────────────────────────────────


def chunk_text(text: str, chunk_size: int = 1024, overlap: int = 128) -> list[str]:
    """Split text into overlapping chunks using paragraph/sentence boundaries.

    Tries to split on paragraph breaks first, then sentences, then words,
    falling back to character-level splits. Adapted from the chunking strategy
    in validatedpatterns-sandbox/vector-embedder.
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return [chunk for chunk in splitter.split_text(text) if chunk.strip()]


def extract_html_text(html: str) -> str:
    """Extract text from an HTML string, stripping scripts/styles/nav."""
    parser = _make_html_parser()
    parser.feed(html)
    return "\n\n".join(parser.parts)
